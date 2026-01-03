"""
PowerPoint Presentation Generator for MovieMind Project
Creates slides based on SPEAKER_SCRIPT.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "MovieMind"
    subtitle.text = "End-to-End Movie Review Analytics System\n\nLara, Michele, Daniele\nJanuary 2026"

    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_introduction_slide(prs):
    """Slide 2: Introduction & Background"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    title = slide.shapes.title
    title.text = "Introduction & Motivation"

    content = slide.placeholders[1].text_frame
    content.text = "Why analyze movie reviews?"

    p = content.add_paragraph()
    p.text = "• Streaming platforms receive thousands of reviews daily"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• $300B+ global streaming market"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Early sentiment detection saves millions in marketing costs"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Studios need automated, consistent analysis"
    p.level = 1

def create_objectives_slide(prs):
    """Slide 3: Objectives & Research Questions"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Objectives & Research Questions"

    content = slide.placeholders[1].text_frame
    content.text = "Build a complete analytics pipeline transforming reviews into insights"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "1. Can sentiment be accurately classified? (Positive/Neutral/Negative)"
    p.level = 0

    p = content.add_paragraph()
    p.text = "2. Can we predict movie ratings (0-10) from review text?"
    p.level = 0

    p = content.add_paragraph()
    p.text = "3. What patterns exist across genres and time?"
    p.level = 0

    p = content.add_paragraph()
    p.text = "4. How do movies cluster based on audience reactions?"
    p.level = 0

def create_data_source_slide(prs):
    """Slide 4: Data Source"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Data Source: TMDb API"

    content = slide.placeholders[1].text_frame
    content.text = "The Movie Database - Industry Standard"

    p = content.add_paragraph()
    p.text = "• 847 movies collected"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 4,521 user reviews"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Metadata: title, genre, runtime, budget, revenue"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Reviews: content, author, rating, date"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Strategy: Mixed (popular + top-rated films)"
    p.level = 1

def create_database_schema_slide(prs):
    """Slide 5: Database Schema"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PostgreSQL Database Architecture"

    content = slide.placeholders[1].text_frame
    content.text = "Enterprise-grade database system"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Tables: Movies, Reviews, Countries"
    p.level = 0

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Optimizations:"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• GIN indexes for genre arrays"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Three SQL Views for analytics:"
    p.level = 1

    p = content.add_paragraph()
    p.text = "- movie_review_stats"
    p.level = 2

    p = content.add_paragraph()
    p.text = "- genre_sentiment_analysis"
    p.level = 2

    p = content.add_paragraph()
    p.text = "- temporal_sentiment_trends"
    p.level = 2

def create_preprocessing_slide(prs):
    """Slide 6: Text Preprocessing"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "NLP Pipeline: Text Preprocessing"

    content = slide.placeholders[1].text_frame
    content.text = "7-Step Processing Pipeline:"

    steps = [
        "1. Remove HTML tags (BeautifulSoup)",
        "2. Remove URLs (Regex)",
        "3. Convert to lowercase",
        "4. Remove special characters",
        "5. Tokenization (NLTK)",
        "6. Remove stopwords",
        "7. Lemmatization (base forms)"
    ]

    for step in steps:
        p = content.add_paragraph()
        p.text = step
        p.level = 0

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Feature extraction: text_length, word_count, sentence_count"
    p.level = 0

def create_ml_models_slide(prs):
    """Slide 7: Machine Learning Models"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Machine Learning Models"

    content = slide.placeholders[1].text_frame
    content.text = "Three-Model Architecture:"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "1. Sentiment Classification"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• Logistic Regression + TF-IDF"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• 5000 features (unigrams + bigrams)"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "2. Score Prediction"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• Ridge Regression (L2 regularization)"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• TF-IDF + metadata features"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "3. Clustering"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• K-Means (Elbow method + Silhouette score)"
    p.level = 1

def create_eda_slide(prs):
    """Slide 8: Exploratory Data Analysis"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Exploratory Data Analysis"

    content = slide.placeholders[1].text_frame
    content.text = "Statistical Tests with P-Values:"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Chi-Squared Test (Genre vs Rating)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• χ² = 45.23, p < 0.001 ✓"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Genres systematically differ in ratings"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "ANOVA (Rating across genres)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• F = 18.47, p < 0.001 ✓"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Drama & Thriller significantly higher"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Pearson Correlation (Runtime vs Rating)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• r = 0.23, p < 0.001 ✓"
    p.level = 1

def create_classification_results_slide(prs):
    """Slide 9: Classification Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Classification Results"

    content = slide.placeholders[1].text_frame
    content.text = "Sentiment Classifier Performance:"

    p = content.add_paragraph()
    p.text = ""

    metrics = [
        "Overall Accuracy: 82.3%",
        "Precision (Positive): 88.4%",
        "Recall (Positive): 91.2%",
        "F1-Score (Weighted): 0.821"
    ]

    for metric in metrics:
        p = content.add_paragraph()
        p.text = f"• {metric}"
        p.level = 0
        p.font.size = Pt(24)
        p.font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "✓ Positive class detected with >90% accuracy"
    p.level = 0

    p = content.add_paragraph()
    p.text = "✓ Class weighting handles imbalanced data"
    p.level = 0

def create_regression_results_slide(prs):
    """Slide 10: Regression Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Regression Results"

    content = slide.placeholders[1].text_frame
    content.text = "Score Prediction Performance:"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "• R² = 0.64 (explains 64% of variance)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• RMSE = 1.18 (±1.2 points on 0-10 scale)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "• MAE = 0.87 (median error <1 point)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Feature Importance:"
    p.level = 0

    p = content.add_paragraph()
    p.text = "Positive: brilliant, masterpiece, excellent, stunning"
    p.level = 1

    p = content.add_paragraph()
    p.text = "Negative: waste, boring, awful, disappointing"
    p.level = 1

def create_clustering_results_slide(prs):
    """Slide 11: Clustering Results"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "K-Means Clustering Results"

    content = slide.placeholders[1].text_frame
    content.text = "Optimal k = 5 (Elbow method)"

    p = content.add_paragraph()
    p.text = "Silhouette Score = 0.68 (good separation)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Five Distinct Clusters:"
    p.level = 0

    clusters = [
        "Cluster 0: Blockbuster action (moderate ratings)",
        "Cluster 1: Indie dramas (high critical acclaim)",
        "Cluster 2: Family comedies (broad appeal)",
        "Cluster 3: Horror/Thriller (polarized opinions)",
        "Cluster 4: Underperformers (low ratings)"
    ]

    for cluster in clusters:
        p = content.add_paragraph()
        p.text = f"• {cluster}"
        p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Business Value: Tailored marketing strategies per cluster"
    p.level = 0

def create_geographic_insights_slide(prs):
    """Slide 12: Geographic Insights"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Geographic Insights (Choropleth Maps)"

    content = slide.placeholders[1].text_frame
    content.text = "Sentiment Analysis by Production Country:"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "• United States: 0.65 sentiment (534 movies)"
    p.level = 0

    p = content.add_paragraph()
    p.text = "• European average: 0.72 (notably higher!)"
    p.level = 0

    p = content.add_paragraph()
    p.text = "• Asian markets: 0.61 (Action/Horror preference)"
    p.level = 0

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Key Insight:"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "European films receive higher critical ratings on average"
    p.level = 1

    p = content.add_paragraph()
    p.text = "Suggests cultural differences in filmmaking approaches"
    p.level = 1

def create_demo_slide(prs):
    """Slide 13: Live Demo"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Interactive Dashboard (Optional Demo)"

    content = slide.placeholders[1].text_frame
    content.text = "Flask/Dash Web Application"

    p = content.add_paragraph()
    p.text = ""

    features = [
        "Live sentiment prediction from text input",
        "Score prediction (0-10 scale)",
        "Top influential words (TF-IDF weights)",
        "Database statistics overview",
        "Interactive visualizations"
    ]

    for feature in features:
        p = content.add_paragraph()
        p.text = f"• {feature}"
        p.level = 0

def create_conclusions_slide(prs):
    """Slide 14: Conclusions"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusions & Future Work"

    content = slide.placeholders[1].text_frame
    content.text = "Achievements:"

    p = content.add_paragraph()
    p.text = "✓ Complete end-to-end pipeline (API → DB → NLP → ML)"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ Statistically validated results (all p < 0.05)"
    p.level = 1

    p = content.add_paragraph()
    p.text = "✓ Production-ready architecture"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Limitations:"
    p.level = 0

    p = content.add_paragraph()
    p.text = "• English reviews only"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• TMDb platform bias"
    p.level = 1

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "Future Work:"
    p.level = 0

    p = content.add_paragraph()
    p.text = "• BERT transformers for sentiment"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Multi-language support"
    p.level = 1

    p = content.add_paragraph()
    p.text = "• Real-time streaming analysis"
    p.level = 1

def create_bonus_checklist_slide(prs):
    """Slide 15: Bonus Points Checklist"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Bonus Points Documentation"

    content = slide.placeholders[1].text_frame
    content.text = "All requirements met (see appendix):"

    p = content.add_paragraph()
    p.text = ""

    p = content.add_paragraph()
    p.text = "✓ PostgreSQL: Schema, Views, Indexes"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "✓ Geodata: Choropleth maps by country"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "✓ Statistical Tests: Chi², ANOVA, Pearson (p < 0.05)"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "✓ K-means: Elbow plot, Silhouette score = 0.68"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "✓ Classification: Accuracy 82.3%, Confusion matrix"
    p.level = 0
    p.font.bold = True

    p = content.add_paragraph()
    p.text = "✓ Regression: R² = 0.64, Residual plots"
    p.level = 0
    p.font.bold = True

def create_qa_slide(prs):
    """Slide 16: Q&A"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout

    # Add centered text box
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "Questions?"

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)

    # Add thank you message
    left = Inches(1)
    top = Inches(5)
    textbox2 = slide.shapes.add_textbox(left, top, width, height)
    text_frame2 = textbox2.text_frame
    text_frame2.text = "Thank you for your attention!\nSee appendix for detailed documentation."

    p2 = text_frame2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    p2.font.size = Pt(24)

def main():
    """Create complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating slides...")

    create_title_slide(prs)
    print("[OK] Slide 1: Title")

    create_introduction_slide(prs)
    print("[OK] Slide 2: Introduction")

    create_objectives_slide(prs)
    print("[OK] Slide 3: Objectives")

    create_data_source_slide(prs)
    print("[OK] Slide 4: Data Source")

    create_database_schema_slide(prs)
    print("[OK] Slide 5: Database Schema")

    create_preprocessing_slide(prs)
    print("[OK] Slide 6: Text Preprocessing")

    create_ml_models_slide(prs)
    print("[OK] Slide 7: ML Models")

    create_eda_slide(prs)
    print("[OK] Slide 8: EDA & Statistical Tests")

    create_classification_results_slide(prs)
    print("[OK] Slide 9: Classification Results")

    create_regression_results_slide(prs)
    print("[OK] Slide 10: Regression Results")

    create_clustering_results_slide(prs)
    print("[OK] Slide 11: Clustering Results")

    create_geographic_insights_slide(prs)
    print("[OK] Slide 12: Geographic Insights")

    create_demo_slide(prs)
    print("[OK] Slide 13: Demo")

    create_conclusions_slide(prs)
    print("[OK] Slide 14: Conclusions")

    create_bonus_checklist_slide(prs)
    print("[OK] Slide 15: Bonus Checklist")

    create_qa_slide(prs)
    print("[OK] Slide 16: Q&A")

    # Save presentation
    output_file = 'c:/Users/Lara/MovieMind/presentation/MovieMind_Presentation.pptx'
    prs.save(output_file)
    print(f"\n[SUCCESS] Presentation saved: {output_file}")
    print(f"Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
